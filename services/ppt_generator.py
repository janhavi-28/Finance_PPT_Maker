from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io
import base64
from typing import Dict, List, Any, Optional
from datetime import datetime
import os

class ProfessionalPPTGenerator:
    """Generate professional PowerPoint presentations with Gamma AI-level quality"""
    
    def __init__(self):
        self.templates = {
            "corporate_blue": {
                "primary_color": RGBColor(31, 78, 121),      # #1f4e79
                "secondary_color": RGBColor(74, 144, 226),    # #4a90e2
                "accent_color": RGBColor(123, 179, 240),      # #7bb3f0
                "text_color": RGBColor(51, 51, 51),          # #333333
                "background_color": RGBColor(255, 255, 255),  # #ffffff
                "success_color": RGBColor(40, 167, 69),       # #28a745
                "warning_color": RGBColor(255, 193, 7),       # #ffc107
                "danger_color": RGBColor(220, 53, 69)         # #dc3545
            },
            "financial_green": {
                "primary_color": RGBColor(45, 90, 39),        # #2d5a27
                "secondary_color": RGBColor(76, 175, 80),     # #4caf50
                "accent_color": RGBColor(129, 199, 132),      # #81c784
                "text_color": RGBColor(51, 51, 51),          # #333333
                "background_color": RGBColor(255, 255, 255),  # #ffffff
                "success_color": RGBColor(102, 187, 106),     # #66bb6a
                "warning_color": RGBColor(255, 152, 0),       # #ff9800
                "danger_color": RGBColor(244, 67, 54)         # #f44336
            },
            "modern_orange": {
                "primary_color": RGBColor(230, 81, 0),        # #e65100
                "secondary_color": RGBColor(255, 152, 0),     # #ff9800
                "accent_color": RGBColor(255, 183, 77),       # #ffb74d
                "text_color": RGBColor(51, 51, 51),          # #333333
                "background_color": RGBColor(255, 255, 255),  # #ffffff
                "success_color": RGBColor(76, 175, 80),       # #4caf50
                "warning_color": RGBColor(255, 193, 7),       # #ffc107
                "danger_color": RGBColor(244, 67, 54)         # #f44336
            },
            "colorful_modern": {
                "bar_colors": [
                    RGBColor(220, 53, 69),   # Red for D
                    RGBColor(255, 99, 132),  # Pink for C
                    RGBColor(40, 167, 69),   # Green for B
                    RGBColor(255, 193, 7)    # Yellow/Beige for A
                ],
                "letter_color": RGBColor(255, 255, 255),  # White letters
                "title_color": RGBColor(0, 0, 0),  # Black title
                "background_color": RGBColor(255, 255, 255)  # White background
            },
            "modern_white": {
                "primary_color": RGBColor(0, 51, 102),        # #003366
                "secondary_color": RGBColor(173, 216, 230),   # #add8e6
                "accent_color": RGBColor(240, 248, 255),      # #f0f8ff
                "text_color": RGBColor(51, 51, 51),          # #333333
                "background_color": RGBColor(255, 255, 255),  # #ffffff
                "success_color": RGBColor(144, 238, 144),     # #90ee90
                "warning_color": RGBColor(255, 228, 196),     # #ffe4c4
                "danger_color": RGBColor(255, 182, 193)       # #ffb6c1
            },
            "executive_dark": {
                "primary_color": RGBColor(25, 25, 25),        # #191919
                "secondary_color": RGBColor(64, 64, 64),      # #404040
                "accent_color": RGBColor(105, 105, 105),      # #696969
                "text_color": RGBColor(255, 255, 255),        # #ffffff
                "background_color": RGBColor(18, 18, 18),     # #121212
                "success_color": RGBColor(76, 175, 80),       # #4caf50
                "warning_color": RGBColor(255, 193, 7),       # #ffc107
                "danger_color": RGBColor(244, 67, 54)         # #f44336
            }
        }
    
    def create_presentation(self, 
                          content: Dict[str, Any], 
                          template: str = "corporate_blue") -> bytes:
        """Create a complete professional presentation"""
        
        prs = Presentation()
        colors = self.templates[template]
        
        # Set slide size to widescreen (16:9)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Create title slide
        self._create_title_slide(prs, content, colors)
        
        # Create agenda slide
        self._create_agenda_slide(prs, content, colors)
        
        # Create content slides
        for slide_data in content.get("slides", []):
            self._create_content_slide(prs, slide_data, colors)
        
        # Save to bytes
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        
        return ppt_io.getvalue()
    
    def _create_title_slide(self, prs, content: Dict, colors: Dict):
        """Create professional title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        if "colorful_modern" in colors:
            # Custom colorful modern template
            self._create_colorful_title_slide(prs, slide, content, colors)
        else:
            # Original templates
            # Set background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = colors["background_color"]
            
            # Title
            title = slide.shapes.title
            title.text = content.get("presentation_title", "Financial Analysis")
            title_frame = title.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.font.name = "Calibri"
            title_para.font.size = Pt(44)
            title_para.font.color.rgb = colors["primary_color"]
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            # Subtitle
            subtitle = slide.placeholders[1]
            subtitle.text = f"Professional Financial Presentation\n{datetime.now().strftime('%B %Y')}"
            subtitle_frame = subtitle.text_frame
            for para in subtitle_frame.paragraphs:
                para.font.name = "Calibri"
                para.font.size = Pt(24)
                para.font.color.rgb = colors["secondary_color"]
                para.alignment = PP_ALIGN.CENTER
            

    def _create_colorful_title_slide(self, prs, slide, content: Dict, colors: Dict):
        """Create colorful modern title slide with vertical bars and letters"""
        # Clear default layout elements
        for shape in slide.shapes:
            shape.element.getparent().remove(shape.element)
        
        # Set white background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["background_color"]
        
        # Create 4 vertical bars on left (each 1/4 of slide width, full height)
        bar_width = Inches(3.33)  # 1/4 of 13.33 inches
        bar_height = Inches(7.5)  # Full height
        bar_x_start = Inches(0)
        
        letters = ['D', 'C', 'B', 'A']
        bar_colors = colors["bar_colors"]
        
        for i, (letter, bar_color) in enumerate(zip(letters, bar_colors)):
            # Bar position
            bar_x = bar_x_start + (i * bar_width)
            
            # Create bar rectangle
            bar_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                bar_x, Inches(0), bar_width, bar_height
            )
            bar_shape.fill.solid()
            bar_shape.fill.fore_color.rgb = bar_color
            bar_shape.line.visible = False  # No border
            
            # Add letter overlay (centered vertically and horizontally in bar)
            letter_shape = slide.shapes.add_textbox(
                bar_x + Inches(0.5), Inches(3), Inches(2.33), Inches(1.5)
            )
            letter_frame = letter_shape.text_frame
            letter_frame.clear()
            letter_para = letter_frame.add_paragraph()
            letter_para.text = letter
            letter_para.font.name = "Arial Black"
            letter_para.font.size = Pt(48)
            letter_para.font.color.rgb = colors["letter_color"]
            letter_para.font.bold = True
            letter_para.alignment = PP_ALIGN.CENTER
            letter_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Title on right side
        title_x = Inches(4)
        title_shape = slide.shapes.add_textbox(
            title_x, Inches(1.5), Inches(9), Inches(2)
        )
        title_frame = title_shape.text_frame
        title_para = title_frame.add_paragraph()
        title_para.text = content.get("presentation_title", "Financial Analysis")
        title_para.font.name = "Calibri"
        title_para.font.size = Pt(44)
        title_para.font.color.rgb = colors["title_color"]
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.LEFT
        
        # Subtitle below title
        subtitle_shape = slide.shapes.add_textbox(
            title_x, Inches(3.5), Inches(9), Inches(1.5)
        )
        subtitle_frame = subtitle_shape.text_frame
        subtitle_para = subtitle_frame.add_paragraph()
        subtitle_para.text = f"Professional Financial Presentation\n{datetime.now().strftime('%B %Y')}"
        subtitle_para.font.name = "Calibri"
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = colors["title_color"]
        subtitle_para.alignment = PP_ALIGN.LEFT
    
    def _create_agenda_slide(self, prs, content: Dict, colors: Dict):
        """Create agenda/outline slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["background_color"]
        
        # Title
        title = slide.shapes.title
        title.text = "Agenda"
        self._format_slide_title(title, colors)
        
        # Content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        # Add agenda items from slides
        slides_data = content.get("slides", [])
        for i, slide_data in enumerate(slides_data[:8]):  # Limit to 8 items
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"{i+1}. {slide_data.get('title', f'Section {i+1}')}"
            p.font.name = "Calibri"
            p.font.size = Pt(20)
            p.font.color.rgb = colors["text_color"]
            p.level = 0
        
        # Add executive summary - repositioned to avoid overlap
        if content.get("executive_summary"):
            summary_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(5), Inches(6), Inches(2)
            )
            summary_frame = summary_box.text_frame
            summary_para = summary_frame.paragraphs[0]
            summary_para.text = "Executive Summary"
            summary_para.font.name = "Calibri"
            summary_para.font.size = Pt(16)
            summary_para.font.color.rgb = colors["primary_color"]
            summary_para.font.bold = True
            
            summary_content = summary_frame.add_paragraph()
            summary_content.text = content["executive_summary"][:100] + "..."  # Truncate to fit
            summary_content.font.name = "Calibri"
            summary_content.font.size = Pt(12)
            summary_content.font.color.rgb = colors["text_color"]
    
    def _create_content_slide(self, prs, slide_data: Dict, colors: Dict):
        """Create individual content slide"""
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["background_color"]
        
        # Title
        title = slide.shapes.title
        title.text = slide_data.get("title", "Content Slide")
        self._format_slide_title(title, colors)
        
        # Main content area - wider left side for better text visibility
        content_left = Inches(0.3)
        content_top = Inches(1.3)
        content_width = Inches(7.5)  # Wider for text dominance
        content_height = Inches(5.8)  # Taller
        
        content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        text_frame = content_box.text_frame
        text_frame.clear()
        text_frame.word_wrap = True  # Ensure wrapping
        
        # Add bullet points - strip existing bullets, no level for clean bullets
        content_items = slide_data.get("content", [])
        for i, item in enumerate(content_items):
            # Thorough stripping of bullets/symbols
            clean_item = item.lstrip("• ").lstrip("•").lstrip("- ").lstrip("-").lstrip("* ").lstrip("*").strip()
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"• {clean_item}"  # Manually add single bullet
            p.font.name = "Calibri"
            p.font.size = Pt(18)  # Slightly smaller for fit
            p.font.color.rgb = colors["text_color"]
            p.alignment = PP_ALIGN.LEFT
            p.level = 0  # No auto-bullets
        
        # Add data points if available - as sub-items without double bullets
        data_points = slide_data.get("data_points", [])
        if data_points:
            for data_point in data_points:
                clean_dp = data_point.lstrip("• ").lstrip("•").lstrip("- ").lstrip("-").strip()
                p = text_frame.add_paragraph()
                p.text = f"  • {clean_dp}"  # Indented sub-bullet
                p.font.name = "Calibri"
                p.font.size = Pt(16)
                p.font.color.rgb = colors["text_color"]  # Same color for consistency
                p.font.bold = True
                p.alignment = PP_ALIGN.LEFT
                p.level = 0
        
        # Add image if available - smaller, right side, lower position
        visual_suggestion = slide_data.get("visual_suggestion", {})
        if visual_suggestion and visual_suggestion.get("image_url"):
            self._add_image_to_slide(slide, visual_suggestion, colors)
        
        # No key insight box for clean slides
    
    def _add_image_to_slide(self, slide, visual_suggestion: Dict, colors: Dict):
        """Add image to slide - smaller, right side, no overlap with text"""
        image_url = visual_suggestion.get("image_url")
        if not image_url:
            return
        
        try:
            # Download image
            import requests
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            
            # Position image in bottom-right corner, slightly bigger size
            image_left = Inches(9.0)  # Adjusted left for bigger image
            image_top = Inches(4.2)  # Adjusted top
            image_width = Inches(3.0)  # Slightly bigger width
            image_height = Inches(2.5)  # Proportional height
            
            # Add image
            img_placeholder = slide.shapes.add_picture(
                io.BytesIO(response.content),
                image_left, image_top, 
                width=image_width, height=image_height
            )
            
            # Add subtle image caption below - smaller, no bold, gray color
            caption_top = image_top + image_height + Inches(0.05)
            caption_shape = slide.shapes.add_textbox(image_left, caption_top, image_width, Inches(0.3))
            caption_frame = caption_shape.text_frame
            caption_para = caption_frame.add_paragraph()
            caption_para.text = visual_suggestion.get("image_description", "Visual Representation")
            caption_para.font.name = "Calibri"
            caption_para.font.size = Pt(10)
            caption_para.font.color.rgb = RGBColor(128, 128, 128)  # Gray for subtlety
            caption_para.font.bold = False
            caption_para.alignment = PP_ALIGN.CENTER
                
        except Exception as e:
            # Skip image if download fails - no placeholder
            pass
    
    def _create_summary_slide(self, prs, content: Dict, colors: Dict):
        """Create summary/recommendations slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["background_color"]
        
        # Title
        title = slide.shapes.title
        title.text = "Key Recommendations"
        self._format_slide_title(title, colors)
        
        # Recommendations
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        recommendations = content.get("key_recommendations", [
            "Strategic recommendation based on analysis",
            "Operational improvement opportunity", 
            "Risk mitigation strategy"
        ])
        
        for i, rec in enumerate(recommendations):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"{i+1}. {rec}"
            p.font.name = "Calibri"
            p.font.size = Pt(20)
            p.font.color.rgb = colors["text_color"]
            p.font.bold = True
            p.level = 0
        
        # Add next steps box
        next_steps_box = slide.shapes.add_textbox(
            Inches(8), Inches(2), Inches(4.5), Inches(4)
        )
        next_steps_frame = next_steps_box.text_frame
        next_steps_para = next_steps_frame.paragraphs[0]
        next_steps_para.text = "Next Steps"
        next_steps_para.font.name = "Calibri"
        next_steps_para.font.size = Pt(18)
        next_steps_para.font.color.rgb = colors["primary_color"]
        next_steps_para.font.bold = True
        
        steps = [
            "Review and validate findings",
            "Implement priority recommendations",
            "Monitor key performance indicators",
            "Schedule follow-up analysis"
        ]
        
        for step in steps:
            step_para = next_steps_frame.add_paragraph()
            step_para.text = f"• {step}"
            step_para.font.name = "Calibri"
            step_para.font.size = Pt(14)
            step_para.font.color.rgb = colors["text_color"]
    
    def _create_appendix_slide(self, prs, content: Dict, colors: Dict):
        """Create appendix slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = colors["background_color"]
        
        # Title
        title = slide.shapes.title
        title.text = "Appendix & Supporting Materials"
        self._format_slide_title(title, colors)
        
        # Appendix content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        appendix_items = content.get("appendix_suggestions", [
            "Detailed financial models and assumptions",
            "Supporting market research and analysis",
            "Risk assessment matrices and scenarios",
            "Additional data sources and methodology"
        ])
        
        for i, item in enumerate(appendix_items):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.name = "Calibri"
            p.font.size = Pt(18)
            p.font.color.rgb = colors["text_color"]
            p.level = 0
        
        # Add contact information
        contact_box = slide.shapes.add_textbox(
            Inches(8), Inches(4), Inches(4.5), Inches(2)
        )
        contact_frame = contact_box.text_frame
        contact_para = contact_frame.paragraphs[0]
        contact_para.text = "Questions & Discussion"
        contact_para.font.name = "Calibri"
        contact_para.font.size = Pt(18)
        contact_para.font.color.rgb = colors["primary_color"]
        contact_para.font.bold = True
        
        contact_content = contact_frame.add_paragraph()
        contact_content.text = "Generated by FinancePPT AI\nProfessional Financial Analysis Platform"
        contact_content.font.name = "Calibri"
        contact_content.font.size = Pt(14)
        contact_content.font.color.rgb = colors["text_color"]
    
    def _format_slide_title(self, title_shape, colors: Dict):
        """Apply consistent title formatting"""
        title_frame = title_shape.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "Calibri"
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = colors["primary_color"]
        title_para.font.bold = True
    
    def save_presentation(self, ppt_bytes: bytes, filename: str) -> str:
        """Save presentation to output directory"""
        output_path = os.path.join("output", filename)
        os.makedirs("output", exist_ok=True)
        
        with open(output_path, "wb") as f:
            f.write(ppt_bytes)
        
        return output_path

# Global PPT generator instance
ppt_generator = ProfessionalPPTGenerator()
